from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create a new presentation
prs = Presentation()


# Helper function to add a slide with title and content
def add_slide(prs, layout_idx, title_text, content_text="", table_data=None):
    slide_layout = prs.slide_layouts[layout_idx]  # 0: Title, 1: Title+Content
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = title_text
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(128, 0, 128)  # Purple
    title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    if content_text:
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.text = content_text
        for p in tf.paragraphs:
            p.font.size = Pt(18)

    if table_data:
        rows, cols = len(table_data), len(table_data[0])
        left = Inches(0.5)
        top = Inches(2)
        width = Inches(8)
        height = Inches(1.5 * rows / 3)  # Rough scaling
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        for r in range(rows):
            for c in range(cols):
                table.cell(r, c).text = str(table_data[r][c])
                table.cell(r, c).text_frame.paragraphs[0].font.size = Pt(14)

    return slide


# Slide 1: Title Slide
add_slide(prs, 0,
          "CS 188: Artificial Intelligence\nLocal and Greedy Search Algorithms\nHill Climbing • Beam Search • Greedy Best-First Search\n\nFall 2025\n[Your Name]")

# Slide 2: Overview
table_data = [
    ["Algorithm", "Complete?", "Optimal?", "Memory", "Real-world use"],
    ["Hill Climbing", "No", "No", "O(1)", "Neural net training, 8-queens"],
    ["Beam Search", "No", "No", "O(beam)", "Machine translation, ChatGPT decode"],
    ["Greedy Best-First", "No", "No", "O(b^d)", "Pathfinding with good heuristic"]
]
add_slide(prs, 1, "Today’s Algorithms",
          "All three are FAST but sacrifice guarantees\nThey trade optimality/completeness for speed",
          table_data=table_data)

# Slide 3: Hill Climbing (1/3)
add_slide(prs, 1, "Hill Climbing = “Greedy Local Search”",
          "Idea:\nFrom current state → look at all neighbors → move to the BEST one (highest value / lowest cost)\nRepeat until no neighbor is better → STOP\n\nLike climbing a mountain in thick fog: always step uphill\n\n[Placeholder: Robot climbing hill with fog]")

# Slide 4: Hill Climbing (2/3)
add_slide(prs, 1, "Hill Climbing Problems",
          "- Local maximum\n- Plateau (flat area)\n- Ridge (slow oscillation)\n\nFixes later → Simulated Annealing, Random Restarts\n\n[Placeholder: Three classic landscapes with robot stuck]")

# Slide 5: Hill Climbing (3/3)
add_slide(prs, 1, "Hill Climbing Pseudocode",
          """current ← random/initial state
loop:
    neighbor ← best successor of current
    if value(neighbor) ≤ value(current): return current
    current ← neighbor

Extremely simple, O(1) memory  
Great when “good enough” is enough!""")

# Slide 6: Greedy Best-First Search
add_slide(prs, 1, "Greedy Best-First Search",
          "Uses heuristic h(n) = estimated distance to goal\nAlways expand the node that looks closest to the goal\n\nPriority queue ordered only by h(n)\n→ Ignores path cost g(n) completely\n\nFast, but can go badly wrong if heuristic misleads!\n\n[Placeholder: Robot with GPS looking at straight-line distance]")

# Slide 7: Greedy Best-First Example
add_slide(prs, 1, "Romania: Arad → Bucharest\nHeuristic = straight-line distance",
          "Path found:\nArad → Sibiu → Făgăraș → Bucharest (379 km)\nOptimal is actually 418 km? Wait — no, greedy got lucky this time!\n\nBut if there’s a tempting dead-end with low h(n) → greedy dives in!\n\n[Placeholder: Map with straight-line arrows]")

# Slide 8: Beam Search (1/3)
add_slide(prs, 1, "Beam Search – Best of Both Worlds",
          "Keep ONLY the k best partial paths at each level (k = beam width)\n\nAt each step:\n- Expand all successors of current beam\n- Score them with heuristic / model score\n- Keep top k, discard everything else forever\n\nbeam width = 1 → pure greedy\nbeam width = ∞ → breadth-first\n\n[Placeholder: Team of k robots walking together]")

# Slide 9: Beam Search (2/3)
add_slide(prs, 1, "Beam Search in Language Generation",
          "Task: generate next word\nModel gives probabilities\n\nBeam width = 4 example:\n\nStep 1 best: The cat\nStep 2 expands to:\nThe cat sat …\nThe cat is …\nThe cat jumped …\nThe cat sleeps …\n\nKeep top 4 → discard millions of others\n\nUsed in Google Translate, GPT decoding, speech recognition\n\n[Placeholder: Robot choir singing four best sentences]")

# Slide 10: Beam Search (3/3)
table_data = [
    ["Beam width", "Behavior", "Quality vs Speed"],
    ["1", "Pure greedy", "Fast, often bad"],
    ["4–32", "Industry sweet spot", "Great balance"],
    ["1000+", "High quality, very slow", "Research"],
    ["∞", "Exact search (BFS)", "Complete"]
]
add_slide(prs, 1, "Beam Search Summary", table_data=table_data)

# Slide 11: Final Comparison
table_data = [
    ["Algorithm", "Memory", "Complete", "Optimal", "When to use"],
    ["Hill Climbing", "O(1)", "No", "No", "Huge spaces, continuous, “good enough” fast"],
    ["Greedy Best-First", "O(b^d)", "No", "No", "Pathfinding with decent heuristic"],
    ["Beam Search", "O(k×depth)", "No", "No", "Sequence generation (text, speech, planning)"],
    ["A* (next lecture)", "O(b^d)", "Yes", "Yes", "When you need the true best path"]
]
add_slide(prs, 1, "Final Cheat Sheet", table_data=table_data)

# Slide 12: Takeaway
add_slide(prs, 1, "One-sentence summary",
          "Hill Climbing = always uphill, no memory\nGreedy Best-First = always toward the goal (as the crow flies)\nBeam Search = keep k best possibilities → powers modern NLP\n\nAll three are incomplete and non-optimal on purpose — because in real life, speed usually beats perfection!\n\n[Placeholder: Robot high-fiving a beam of light]\n\nQuestions?")

# Save the file
prs.save('AI_Local_Greedy_Search.pptx')
print("PPTX file created successfully!")